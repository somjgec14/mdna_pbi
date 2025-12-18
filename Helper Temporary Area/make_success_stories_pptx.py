from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

#Pallete
PINE = RGBColor(30, 107, 78) # pine green
HOLLY = RGBColor(198, 40, 40) # holly red
GOLD = RGBColor(212, 175, 55) # gold
TEAL = RGBColor(27, 127, 158) # teal
ORANGE = RGBColor(242, 153, 74) # warm orange
BLUEDK = RGBColor(14, 58, 87) # deep accent
WHITE = RGBColor(255, 255, 255)
GRAYLT = RGBColor(230, 236, 242)
palette = [TEAL, PINE, ORANGE, HOLLY, GOLD]

items = [
"RACE OWIT migration successful with additional useful features; appreciated by business. A critical prerequisite for PBF shutdown.",
"ECMFuture UserCockpit Go‑Live—helping users with a one‑stop overview of engineering change requests in their responsibility.",
"Project Risk Status Forecast—75–83% accuracy predicting risk.",
"Text‑to‑Query PoC.",
"MCR/CSS EY Audit 2025.",
"Efficient platform management using internally developed tools.",
"Initiated cost saving of 1 m€: completed transition/migration of 14 applications out of PBF to enable PBF shutdown.",
"Pay‑per‑activation process established with Daimler Truck.",
"All customer feedback from CX5 ratings scoring 5 out of 5 stars.",
"More than 10 plants piloted.",
"High customer satisfaction with promoter users.",
"Onsite visit in China for direct demonstration and requirements collection.",
"Very good customer satisfaction and positive feedback leading to new use cases for next year.",
"Successfully turned the PoC into an MVP.",
"Team workshops across various locations in India, Germany, and Vietnam.",
"Second place in X‑Perience Award.",
"Click Dashboard—Organization Views.",
]

slide1_items = items[:9]
slide2_items = items[9:]

prs = Presentation()

prs.slide_width = Inches(13.333) # 1280 x 720 pt equivalent
prs.slide_height = Inches(7.5)

title_layout = prs.slide_layouts[5] # blank layout (for full control)

def add_title_and_decor(slide, title_text):
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(10.5), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text + " 🎄✨🎁"
    p.alignment = PP_ALIGN.LEFT
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = BLUEDK
    # Light “snowflake” corner accents
    for (x, y) in [(0.3, 0.2), (12.7, 0.2)]:
        flake = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(0.6), Inches(0.6))
        tff = flake.text_frame
        tff.clear()
        rp = tff.paragraphs[0]
        r = rp.add_run()
        r.text = "❄"
        r.font.size = Pt(22)
        r.font.color.rgb = GRAYLT
    # Logo placeholder (top right)
    logo = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.3), Inches(0.25), Inches(1.7), Inches(0.8))
    logo.fill.solid()
    logo.fill.fore_color.rgb = RGBColor(255, 255, 255)
    logo.line.color.rgb = GRAYLT
    logo.line.width = Pt(1.5)
    logo.adjustments[0] = 0.3
    t = logo.text_frame
    t.text = "Logo"
    t.paragraphs[0].alignment = PP_ALIGN.CENTER
    t.paragraphs[0].runs[0].font.size = Pt(16)
    t.paragraphs[0].runs[0].font.color.rgb = BLUEDK

def add_pill(slide, left, top, width, height, text, fill_rgb):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = RGBColor(255, 255, 255)
    shape.line.width = Pt(2)
    try:
        shape.shadow.inherit = False
        shape.shadow.visible = True
        shape.shadow.blur_radius = Pt(6)
        shape.shadow.distance = Pt(3)
        shape.shadow.angle = 2700000 # 270 degrees in OOXML units
        shape.shadow.transparency = 0.3
    except Exception:
    pass
    shape.adjustments[0] = 0.3 # more rounded
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(18)
    r.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)

def layout_boxes(slide, items_list):
    # Grid: 3 columns x 3 rows (or 3 x 3/2) with comfortable spacing
    cols = 3
    rows = 3
    # Compute needed rows
    rows = (len(items_list) + cols - 1) // cols
    left0 = Inches(0.8)
    top0 = Inches(1.3)
    gapx = Inches(0.4)
    gapy = Inches(0.5)
    box_w = (prs.slide_width - left02 - gapx(cols-1))
    box_w = box_w / cols
    box_h = Inches(1.4)
    for i, text in enumerate(items_list):
        r = i // cols
        c = i % cols
        left = left0 + c * (box_w + gapx)
        top  = top0 + r * (box_h + gapy)
        color = palette[i % len(palette)]
        add_pill(slide, left, top, box_w, box_h, text, color)

s1 = prs.slides.add_slide(title_layout)
add_title_and_decor(s1, "Success Stories")
layout_boxes(s1, slide1_items)

s2 = prs.slides.add_slide(title_layout)
add_title_and_decor(s2, "Success Stories (cont.)")
layout_boxes(s2, slide2_items)

out_name = "Success_Stories_Christmas.pptx"
prs.save(out_name)
print(f"Created {out_name}")
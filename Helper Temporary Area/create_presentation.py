# Filename: create_presentation.py
from pptx import Presentation
from pptx.util import Inches, Pt

# --- Create Presentation ---
prs = Presentation()

# --- Slide 1: Title Slide ---
slide_layout = prs.slide_layouts[0] # 0 is the title slide layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Project Phoenix: Q3 Kick-off"
subtitle.text = "A Strategic Overview for the Coming Quarter\nPresented by: The Core Team"
# --- Slide 2: Agenda ---
slide_layout = prs.slide_layouts[1] # 1 is 'Title and Content'
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Agenda"
content_shape = slide.shapes.placeholders[1]
tf = content_shape.text_frame
tf.clear() # Clear existing bullet points

p1 = tf.add_paragraph()
p1.text = "Review of Q2 Performance"
p1.level = 0

p2 = tf.add_paragraph()
p2.text = "Introduction to Project Phoenix"
p2.level = 0

p3 = tf.add_paragraph()
p3.text = "Key Objectives & Milestones for Q3"
p3.level = 0

p4 = tf.add_paragraph()
p4.text = "Resource Allocation and Team Roles"
p4.level = 0

p5 = tf.add_paragraph()
p5.text = "Next Steps & Q&A"
p5.level = 0

# --- Slide 3: Content with an Image ---
slide_layout = prs.slide_layouts[5] # 5 is 'Title Only'
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Our Primary Goal: Market Expansion"

# Add an image (replace 'image.png' with your image file)
try:
    img_path = 'image.png'
    # Add an image placeholder if you don't have one
    with open(img_path, 'wb') as f:
        f.write(b'\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x10\x00\x00\x00\x10\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\x01sRGB\x00\xae\xce\x1c\xe9\x00\x00\x00\tpHYs\x00\x00\x0b\x13\x00\x00\x0b\x13\x01\x00\x9a\x9c\x18\x00\x00\x00\x07tIME\x07\xe4\x05\x0f\x0e\x0f\x3a\x9e\xd7U\x00\x00\x00\x19tEXtSoftware\x00www.inkscape.org\x9b\xee<\x1a\x00\x00\x00-IDAT\x08\xd7c\xf8\xff\xff?\x03\x0c\x00\x0c\xc3\x01\x07\x00\xf3\xbf\xfb\xef\x0f8\x10\x00\x00\x00\x00\x00\x00IEND\xaeB`\x82') # A simple 16x16 px red dot
    left = Inches(1.5)
    top = Inches(2.0)
    height = Inches(4.0)
    pic = slide.shapes.add_picture(img_path, left, top, height=height)
except FileNotFoundError:
    print("Warning: 'image.png' not found. Skipping image on slide 3.")

# Add a text box
txBox = slide.shapes.add_textbox(Inches(5.5), Inches(2.5), Inches(4), Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "Targeting a 15% increase in market share by EOY."
p.font.bold = True
p.font.size = Pt(24)


# --- Slide 4: Data Table ---
slide_layout = prs.slide_layouts[1] # 'Title and Content'
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Key Performance Indicators (KPIs)"

rows, cols = 4, 3
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(2.0)
table.columns[2].width = Inches(2.0)

# Write table headers
table.cell(0, 0).text = 'Metric'
table.cell(0, 1).text = 'Q2 Actual'
table.cell(0, 2).text = 'Q3 Target'

# Write data
table.cell(1, 0).text = 'User Acquisition'
table.cell(1, 1).text = '1.2M'
table.cell(1, 2).text = '1.5M'

table.cell(2, 0).text = 'Customer Churn'
table.cell(2, 1).text = '5.4%'
table.cell(2, 2).text = '< 4.0%'

table.cell(3, 0).text = 'Revenue (USD)'
table.cell(3, 1).text = '$450K'
table.cell(3, 2).text = '$550K'

# --- Slide 5: Thank You / Q&A ---
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Thank You"
subtitle.text = "Questions & Answers"

# --- Save the Presentation ---
prs.save("Project_Phoenix_Kickoff.pptx")

print("Presentation 'Project_Phoenix_Kickoff.pptx' created successfully!")